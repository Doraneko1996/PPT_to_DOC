import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT

def setup_document():
    """Thiết lập định dạng tài liệu"""
    doc = Document()
    
    # Thiết lập kích thước trang A4 và lề
    section = doc.sections[0]
    section.page_height = Inches(11.69)  # A4 height
    section.page_width = Inches(8.27)    # A4 width
    section.left_margin = Inches(1.18)   # 3 cm
    section.right_margin = Inches(0.79)  # 2 cm
    section.top_margin = Inches(0.79)    # 2 cm
    section.bottom_margin = Inches(0.79) # 2 cm
    
    return doc

def copy_text_format(ppt_run, doc_run, is_title=False):
    """Sao chép và định dạng văn bản"""
    # Thiết lập font Times New Roman và cỡ chữ
    doc_run.font.name = 'Times New Roman'
    doc_run.font.size = Pt(14 if is_title else 12)
    
    # Copy định dạng bold, italic, underline
    doc_run.font.bold = ppt_run.font.bold
    doc_run.font.italic = ppt_run.font.italic
    doc_run.font.underline = ppt_run.font.underline

def copy_table(shape, doc):
    """Sao chép bảng từ PowerPoint sang Word"""
    if not hasattr(shape, "table"):
        return
    
    ppt_table = shape.table
    rows = len(ppt_table.rows)
    cols = len(ppt_table.columns)
    
    # Tạo bảng trong Word
    doc_table = doc.add_table(rows=rows, cols=cols)
    doc_table.style = 'Table Grid'
    
    # Copy nội dung và định dạng từng cell
    for i in range(rows):
        for j in range(cols):
            ppt_cell = ppt_table.cell(i, j)
            doc_cell = doc_table.cell(i, j)
            
            for paragraph in ppt_cell.text_frame.paragraphs:
                doc_paragraph = doc_cell.paragraphs[0]
                for run in paragraph.runs:
                    doc_run = doc_paragraph.add_run(run.text)
                    copy_text_format(run, doc_run)
    
    # Thêm khoảng trống sau bảng
    doc.add_paragraph()

def extract_smartart_text(shape, doc):
    """Trích xuất văn bản từ SmartArt"""
    try:
        # Kiểm tra nếu là SmartArt
        if hasattr(shape, 'graphic'):
            # Lấy tất cả các node text trong SmartArt
            for node in shape._element.iter():
                # Tìm tất cả các thẻ <a:t> chứa text
                if node.tag.endswith('}t'):
                    text = node.text
                    if text and text.strip():
                        # Bỏ qua các text không mong muốn
                        if not any(skip in text for skip in ['Tuần', 'Tiết']):
                            paragraph = doc.add_paragraph()
                            run = paragraph.add_run(text.strip())
                            run.font.name = 'Times New Roman'
                            run.font.size = Pt(14)
                            
                            # Thêm dấu chấm câu nếu cần
                            if not text.strip().endswith(('.', ':', '?', '!')):
                                run = paragraph.add_run('.')
                                run.font.name = 'Times New Roman'
                                run.font.size = Pt(14)
    except Exception as e:
        # Ghi log chi tiết hơn để debug
        print(f"Chi tiết lỗi khi xử lý SmartArt: {str(e)}")

def should_skip_text(text, is_first_slide):
    """Kiểm tra xem có nên bỏ qua text này không"""
    if is_first_slide:
        # Bỏ qua các dòng có chứa "Tuần" hoặc "Tiết" ở slide đầu
        skip_patterns = ["Tuần", "Tiết"]
        return any(pattern in text for pattern in skip_patterns)
    return False

def process_group_shapes(group_shape, doc, is_first_slide=False):
    """Xử lý các shape trong group"""
    if hasattr(group_shape, 'shapes'):
        for shape in group_shape.shapes:
            # Đệ quy nếu gặp group con
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                process_group_shapes(shape, doc, is_first_slide)
            
            # Xử lý text frame trong group
            elif hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and not should_skip_text(text, is_first_slide):
                        doc_paragraph = doc.add_paragraph()
                        
                        # Kiểm tra tiêu đề
                        is_title = any(run.font.size and run.font.size.pt >= 28 for run in paragraph.runs)
                        
                        if is_title:
                            doc_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        
                        # Copy text và định dạng
                        for run in paragraph.runs:
                            doc_run = doc_paragraph.add_run(run.text)
                            copy_text_format(run, doc_run, is_title)

def ppt_to_word(ppt_path, word_path):
    """Chuyển đổi PowerPoint sang Word"""
    prs = Presentation(ppt_path)
    doc = setup_document()
    
    for slide_index, slide in enumerate(prs.slides):
        is_first_slide = (slide_index == 0)
        
        for shape in slide.shapes:
            try:
                # Xử lý group shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    process_group_shapes(shape, doc, is_first_slide)
                
                # Xử lý văn bản thông thường
                elif shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and not should_skip_text(text, is_first_slide):
                            doc_paragraph = doc.add_paragraph()
                            
                            # Kiểm tra tiêu đề
                            is_title = any(run.font.size and run.font.size.pt >= 28 for run in paragraph.runs)
                            
                            if is_title:
                                doc_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            
                            # Copy text và định dạng
                            for run in paragraph.runs:
                                doc_run = doc_paragraph.add_run(run.text)
                                copy_text_format(run, doc_run, is_title)
                
                # Xử lý bảng
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    copy_table(shape, doc)
                
                # Xử lý SmartArt và các GRAPHIC_FRAME khác
                elif shape.shape_type == MSO_SHAPE_TYPE.GRAPHIC_FRAME:
                    extract_smartart_text(shape, doc)
                    
            except Exception as e:
                # Không in thông báo lỗi GRAPHIC_FRAME nữa
                if "GRAPHIC_FRAME" not in str(e):
                    print(f"Lỗi khi xử lý shape: {str(e)}")
                continue
    
    # Lưu file word
    doc.save(word_path)

def process_all_files():
    """Xử lý tất cả các file trong thư mục input"""
    input_dir = "inputs"
    output_dir = "outputs"
    
    if not os.path.exists(input_dir):
        os.makedirs(input_dir)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    ppt_files = [f for f in os.listdir(input_dir) if f.endswith(('.ppt', '.pptx'))]
    
    if not ppt_files:
        print("Không tìm thấy file PowerPoint nào trong thư mục inputs!")
        return
    
    print(f"Tìm thấy {len(ppt_files)} file PowerPoint cần xử lý.")
    
    for ppt_file in ppt_files:
        try:
            input_path = os.path.join(input_dir, ppt_file)
            output_file = os.path.splitext(ppt_file)[0] + '.docx'
            output_path = os.path.join(output_dir, output_file)
            
            print(f"\nĐang xử lý file: {ppt_file}")
            ppt_to_word(input_path, output_path)
            print(f"Đã chuyển đổi thành công: {output_file}")
            
        except Exception as e:
            print(f"Lỗi khi xử lý file {ppt_file}: {str(e)}")

def main():
    try:
        print("=== Chương trình chuyển đổi PowerPoint sang Word ===")
        print("Đang bắt đầu quá trình chuyển đổi...")
        process_all_files()
        print("\nHoàn thành quá trình chuyển đổi!")
        print("Vui lòng kiểm tra thư mục 'outputs' để xem kết quả.")
        
    except Exception as e:
        print(f"Có lỗi xảy ra trong quá trình thực thi: {str(e)}")
    
    input("\nNhấn Enter để thoát chương trình...")

if __name__ == "__main__":
    main()
